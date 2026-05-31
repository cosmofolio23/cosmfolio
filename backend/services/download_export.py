"""
Download & Export Service
Phase 6: Task 6.3 - PDF, ZIP, self-contained HTML, and PowerPoint exports
"""

import logging
import zipfile
import io
import base64
from typing import Dict, Any, Optional, List, BinaryIO
from enum import Enum
from datetime import datetime
import json
import os

logger = logging.getLogger(__name__)


class ExportFormat(str, Enum):
    """Supported export formats"""
    PDF = "pdf"
    HTML = "html"
    ZIP = "zip"
    POWERPOINT = "powerpoint"
    SELF_CONTAINED_HTML = "self_contained_html"


class DownloadExportService:
    """Service for exporting portfolios in multiple formats"""

    def __init__(self):
        """Initialize download export service"""

        self.formats = [
            ExportFormat.PDF,
            ExportFormat.HTML,
            ExportFormat.ZIP,
            ExportFormat.SELF_CONTAINED_HTML,
        ]

        # File size limits (in bytes)
        self.max_file_size = 100 * 1024 * 1024  # 100MB
        self.max_zip_size = 500 * 1024 * 1024   # 500MB

        logger.info("Download Export Service initialized")

    def export_portfolio_pdf(
        self,
        portfolio_id: str,
        portfolio_data: Dict[str, Any],
        pdf_binary: bytes,
        filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Prepare PDF for download

        Args:
            portfolio_id: Portfolio ID
            portfolio_data: Portfolio data
            pdf_binary: PDF file binary content
            filename: Custom filename (optional)

        Returns:
            Download metadata
        """

        try:
            file_size = len(pdf_binary)

            if file_size > self.max_file_size:
                raise ValueError(f"PDF exceeds maximum size of {self.max_file_size} bytes")

            title = portfolio_data.get("title", "portfolio")
            filename = filename or f"{title.lower().replace(' ', '-')}.pdf"

            export_data = {
                "format": ExportFormat.PDF.value,
                "filename": filename,
                "file_size_bytes": file_size,
                "file_size_mb": round(file_size / (1024 * 1024), 2),
                "mime_type": "application/pdf",
                "portfolio_id": portfolio_id,
                "created_at": datetime.utcnow().isoformat(),
                "binary": pdf_binary,
            }

            logger.info(
                f"PDF export prepared for portfolio {portfolio_id}: "
                f"{filename} ({file_size} bytes)"
            )

            return export_data

        except Exception as e:
            logger.error(f"Error preparing PDF export: {str(e)}")
            raise

    def export_portfolio_html(
        self,
        portfolio_id: str,
        portfolio_data: Dict[str, Any],
        html_content: str,
        filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Prepare HTML for download

        Args:
            portfolio_id: Portfolio ID
            portfolio_data: Portfolio data
            html_content: HTML content string
            filename: Custom filename

        Returns:
            Download metadata
        """

        try:
            html_binary = html_content.encode('utf-8')
            file_size = len(html_binary)

            title = portfolio_data.get("title", "portfolio")
            filename = filename or f"{title.lower().replace(' ', '-')}.html"

            export_data = {
                "format": ExportFormat.HTML.value,
                "filename": filename,
                "file_size_bytes": file_size,
                "file_size_mb": round(file_size / (1024 * 1024), 2),
                "mime_type": "text/html",
                "portfolio_id": portfolio_id,
                "created_at": datetime.utcnow().isoformat(),
                "binary": html_binary,
            }

            logger.info(
                f"HTML export prepared for portfolio {portfolio_id}: "
                f"{filename} ({file_size} bytes)"
            )

            return export_data

        except Exception as e:
            logger.error(f"Error preparing HTML export: {str(e)}")
            raise

    def export_portfolio_self_contained_html(
        self,
        portfolio_id: str,
        portfolio_data: Dict[str, Any],
        html_content: str,
        assets: Dict[str, bytes],
        filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Create self-contained HTML (embed CSS/fonts/images)

        Args:
            portfolio_id: Portfolio ID
            portfolio_data: Portfolio data
            html_content: HTML content
            assets: Dict of asset_name -> binary content
            filename: Custom filename

        Returns:
            Download metadata with embedded content
        """

        try:
            # Embed CSS
            html_with_css = self._embed_css(html_content, assets)

            # Embed fonts
            html_with_fonts = self._embed_fonts(html_with_css, assets)

            # Embed images as base64
            final_html = self._embed_images_base64(html_with_fonts, assets)

            html_binary = final_html.encode('utf-8')
            file_size = len(html_binary)

            if file_size > self.max_file_size:
                raise ValueError(
                    f"Self-contained HTML exceeds maximum size. "
                    f"Try optimizing images or reducing content."
                )

            title = portfolio_data.get("title", "portfolio")
            filename = filename or f"{title.lower().replace(' ', '-')}-standalone.html"

            export_data = {
                "format": ExportFormat.SELF_CONTAINED_HTML.value,
                "filename": filename,
                "file_size_bytes": file_size,
                "file_size_mb": round(file_size / (1024 * 1024), 2),
                "mime_type": "text/html",
                "portfolio_id": portfolio_id,
                "created_at": datetime.utcnow().isoformat(),
                "is_self_contained": True,
                "embedded_assets": {
                    "css_count": len([a for a in assets if a.endswith('.css')]),
                    "font_count": len([a for a in assets if a.endswith(('.woff', '.woff2', '.ttf'))]),
                    "image_count": len([a for a in assets if a.endswith(('.jpg', '.png', '.svg'))]),
                },
                "binary": html_binary,
            }

            logger.info(
                f"Self-contained HTML prepared for portfolio {portfolio_id}: "
                f"{filename} ({file_size} bytes)"
            )

            return export_data

        except Exception as e:
            logger.error(f"Error creating self-contained HTML: {str(e)}")
            raise

    def export_portfolio_zip(
        self,
        portfolio_id: str,
        portfolio_data: Dict[str, Any],
        files: Dict[str, bytes],
        filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Create ZIP export with portfolio + assets + metadata

        Args:
            portfolio_id: Portfolio ID
            portfolio_data: Portfolio data
            files: Dict of filename -> binary content
            filename: Custom ZIP filename

        Returns:
            Download metadata
        """

        try:
            # Create ZIP in memory
            zip_buffer = io.BytesIO()

            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                # Add portfolio files
                for file_name, file_content in files.items():
                    zip_file.writestr(file_name, file_content)

                # Add metadata.json
                metadata = {
                    "portfolio_id": portfolio_id,
                    "title": portfolio_data.get("title", ""),
                    "description": portfolio_data.get("description", ""),
                    "author": portfolio_data.get("author", ""),
                    "exported_at": datetime.utcnow().isoformat(),
                    "format_version": "1.0",
                }

                metadata_json = json.dumps(metadata, indent=2).encode('utf-8')
                zip_file.writestr("metadata.json", metadata_json)

                # Add README
                readme_content = self._generate_zip_readme(portfolio_data)
                zip_file.writestr("README.md", readme_content)

            zip_binary = zip_buffer.getvalue()
            file_size = len(zip_binary)

            if file_size > self.max_zip_size:
                raise ValueError(
                    f"ZIP exceeds maximum size of {self.max_zip_size} bytes"
                )

            title = portfolio_data.get("title", "portfolio")
            filename = filename or f"{title.lower().replace(' ', '-')}-export.zip"

            export_data = {
                "format": ExportFormat.ZIP.value,
                "filename": filename,
                "file_size_bytes": file_size,
                "file_size_mb": round(file_size / (1024 * 1024), 2),
                "mime_type": "application/zip",
                "portfolio_id": portfolio_id,
                "created_at": datetime.utcnow().isoformat(),
                "file_count": len(files) + 2,  # +2 for metadata.json and README.md
                "binary": zip_binary,
            }

            logger.info(
                f"ZIP export created for portfolio {portfolio_id}: "
                f"{filename} ({file_size} bytes) with {len(files)} files"
            )

            return export_data

        except Exception as e:
            logger.error(f"Error creating ZIP export: {str(e)}")
            raise

    def export_portfolio_powerpoint(
        self,
        portfolio_id: str,
        portfolio_data: Dict[str, Any],
        slides_data: List[Dict[str, Any]],
        filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Create PowerPoint presentation export

        Args:
            portfolio_id: Portfolio ID
            portfolio_data: Portfolio data
            slides_data: List of slide data (title, content, images)
            filename: Custom filename

        Returns:
            Download metadata
        """

        try:
            # Try to import python-pptx, fall back if not available
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
            except ImportError:
                logger.warning("python-pptx not installed, falling back to stub")
                return self._create_powerpoint_stub(
                    portfolio_id, portfolio_data, filename
                )

            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = portfolio_data.get("title", "Portfolio")
            subtitle.text = portfolio_data.get("author", "") or "Architecture Portfolio"

            # Add content slides
            for slide_data in slides_data[:20]:  # Limit to 20 slides
                slide = self._add_powerpoint_slide(
                    prs, slide_data, portfolio_data
                )

            # Save to bytes
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_binary = pptx_buffer.getvalue()
            file_size = len(pptx_binary)

            if file_size > self.max_file_size:
                raise ValueError(f"PowerPoint exceeds maximum size")

            title = portfolio_data.get("title", "portfolio")
            filename = filename or f"{title.lower().replace(' ', '-')}.pptx"

            export_data = {
                "format": ExportFormat.POWERPOINT.value,
                "filename": filename,
                "file_size_bytes": file_size,
                "file_size_mb": round(file_size / (1024 * 1024), 2),
                "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "portfolio_id": portfolio_id,
                "created_at": datetime.utcnow().isoformat(),
                "slide_count": len(prs.slides),
                "binary": pptx_binary,
            }

            logger.info(
                f"PowerPoint export created for portfolio {portfolio_id}: "
                f"{filename} ({file_size} bytes) with {len(prs.slides)} slides"
            )

            return export_data

        except Exception as e:
            logger.error(f"Error creating PowerPoint export: {str(e)}")
            raise

    def batch_export(
        self,
        portfolio_ids: List[str],
        portfolio_data_map: Dict[str, Dict[str, Any]],
        format: str,
        files_map: Dict[str, Dict[str, bytes]],
    ) -> Dict[str, Any]:
        """
        Batch export multiple portfolios

        Args:
            portfolio_ids: List of portfolio IDs
            portfolio_data_map: Map of portfolio_id -> portfolio_data
            format: Export format
            files_map: Map of portfolio_id -> files dict

        Returns:
            Batch export metadata
        """

        try:
            results = []
            total_size = 0

            for portfolio_id in portfolio_ids:
                portfolio_data = portfolio_data_map.get(portfolio_id, {})
                files = files_map.get(portfolio_id, {})

                try:
                    if format == ExportFormat.ZIP.value:
                        result = self.export_portfolio_zip(
                            portfolio_id, portfolio_data, files
                        )
                    # Add other formats as needed
                    else:
                        raise ValueError(f"Unsupported format for batch: {format}")

                    results.append({
                        "portfolio_id": portfolio_id,
                        "status": "success",
                        "filename": result["filename"],
                        "file_size": result["file_size_bytes"],
                    })

                    total_size += result["file_size_bytes"]

                except Exception as e:
                    results.append({
                        "portfolio_id": portfolio_id,
                        "status": "failed",
                        "error": str(e),
                    })

            batch_data = {
                "status": "batch_exported",
                "format": format,
                "total_portfolios": len(portfolio_ids),
                "successful": len([r for r in results if r["status"] == "success"]),
                "failed": len([r for r in results if r["status"] == "failed"]),
                "total_size_bytes": total_size,
                "total_size_mb": round(total_size / (1024 * 1024), 2),
                "results": results,
            }

            logger.info(f"Batch export completed: {batch_data['successful']}/{batch_data['total_portfolios']}")

            return batch_data

        except Exception as e:
            logger.error(f"Error in batch export: {str(e)}")
            raise

    def get_export_info(self, format: str) -> Dict[str, Any]:
        """Get info about export format"""

        info = {
            "pdf": {
                "name": "PDF Document",
                "description": "Professional PDF with styling",
                "mime_type": "application/pdf",
                "extension": ".pdf",
                "suitable_for": ["printing", "sharing", "archiving"],
            },
            "html": {
                "name": "HTML File",
                "description": "Interactive HTML for web browsers",
                "mime_type": "text/html",
                "extension": ".html",
                "suitable_for": ["web", "email", "online sharing"],
            },
            "self_contained_html": {
                "name": "Self-Contained HTML",
                "description": "Single file with embedded CSS, fonts, and images",
                "mime_type": "text/html",
                "extension": ".html",
                "suitable_for": ["offline viewing", "email attachment", "no dependencies"],
            },
            "zip": {
                "name": "ZIP Archive",
                "description": "Compressed package with all portfolio files",
                "mime_type": "application/zip",
                "extension": ".zip",
                "suitable_for": ["archiving", "sharing", "backup"],
            },
            "powerpoint": {
                "name": "PowerPoint Presentation",
                "description": "Editable presentation slides",
                "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "extension": ".pptx",
                "suitable_for": ["presentations", "client meetings", "editing"],
            },
        }

        return info.get(format, {})

    # ==================== HELPER METHODS ====================

    def _embed_css(self, html: str, assets: Dict[str, bytes]) -> str:
        """Embed CSS into HTML"""

        result = html
        for asset_name, content in assets.items():
            if asset_name.endswith('.css'):
                css_content = content.decode('utf-8', errors='ignore')
                # Find and replace CSS link with style tag
                tag = f'<link rel="stylesheet" href="{asset_name}">'
                replacement = f'<style>{css_content}</style>'
                result = result.replace(tag, replacement)

        return result

    def _embed_fonts(self, html: str, assets: Dict[str, bytes]) -> str:
        """Embed fonts as base64"""

        result = html
        for asset_name, content in assets.items():
            if asset_name.endswith(('.woff', '.woff2', '.ttf')):
                font_base64 = base64.b64encode(content).decode('utf-8')
                # Create @font-face rule
                ext = asset_name.split('.')[-1]
                font_type = 'woff2' if ext == 'woff2' else 'woff' if ext == 'woff' else 'truetype'

                font_face = f'''
<style>
@font-face {{
    font-family: '{asset_name.split('.')[0]}';
    src: url('data:font/{font_type};base64,{font_base64}') format('{font_type}');
}}
</style>
'''
                result = result.replace(f'url("{asset_name}")', f"url('data:font/{font_type};base64,{font_base64}')")

        return result

    def _embed_images_base64(self, html: str, assets: Dict[str, bytes]) -> str:
        """Embed images as base64 data URLs"""

        result = html
        for asset_name, content in assets.items():
            if asset_name.endswith(('.jpg', '.jpeg', '.png', '.svg')):
                image_base64 = base64.b64encode(content).decode('utf-8')

                # Determine MIME type
                ext = asset_name.split('.')[-1].lower()
                mime_types = {
                    'jpg': 'image/jpeg',
                    'jpeg': 'image/jpeg',
                    'png': 'image/png',
                    'svg': 'image/svg+xml',
                }
                mime_type = mime_types.get(ext, 'image/jpeg')

                data_url = f"data:{mime_type};base64,{image_base64}"
                result = result.replace(f'src="{asset_name}"', f'src="{data_url}"')
                result = result.replace(f"src='{asset_name}'", f"src='{data_url}'")

        return result

    def _generate_zip_readme(self, portfolio_data: Dict[str, Any]) -> str:
        """Generate README for ZIP archive"""

        title = portfolio_data.get("title", "Portfolio")
        author = portfolio_data.get("author", "")
        description = portfolio_data.get("description", "")

        readme = f"""# {title}

**Author:** {author}

## Overview

{description}

## Contents

This ZIP archive contains your complete portfolio export with:

- **index.html** - Main portfolio page
- **styles/** - CSS stylesheets
- **images/** - Portfolio images
- **assets/** - Additional resources
- **metadata.json** - Portfolio metadata
- **README.md** - This file

## How to Use

1. Extract the ZIP archive to a folder
2. Open `index.html` in your web browser
3. View your portfolio with all styling intact

## Requirements

- A modern web browser (Chrome, Firefox, Safari, Edge)
- No server or internet connection required

## Features

✓ Fully responsive design
✓ Self-contained (no external dependencies)
✓ Print-friendly styling
✓ SEO optimized

## Support

For more information, visit CosmoFolio.com

---

*Exported on {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')} UTC*
"""

        return readme

    def _create_powerpoint_stub(
        self,
        portfolio_id: str,
        portfolio_data: Dict[str, Any],
        filename: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Create stub response when python-pptx not installed"""

        logger.warning(
            "PowerPoint export requires python-pptx. "
            "Install with: pip install python-pptx"
        )

        title = portfolio_data.get("title", "portfolio")
        filename = filename or f"{title.lower().replace(' ', '-')}.pptx"

        return {
            "format": ExportFormat.POWERPOINT.value,
            "filename": filename,
            "file_size_bytes": 0,
            "portfolio_id": portfolio_id,
            "status": "unavailable",
            "message": "PowerPoint export requires python-pptx. Install with: pip install python-pptx",
        }

    def _add_powerpoint_slide(
        self,
        prs,
        slide_data: Dict[str, Any],
        portfolio_data: Dict[str, Any],
    ):
        """Add content slide to presentation"""

        try:
            from pptx.util import Inches, Pt
        except ImportError:
            return None

        # Use blank slide layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        if "title" in slide_data:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True

        # Add content
        if "content" in slide_data:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = slide_data["content"]

        return slide


# ==================== SINGLETON INSTANCE ====================

_download_export_service = None

def get_download_export_service() -> DownloadExportService:
    """Get or create download export service singleton"""
    global _download_export_service
    if _download_export_service is None:
        _download_export_service = DownloadExportService()
    return _download_export_service
